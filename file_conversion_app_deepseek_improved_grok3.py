import gradio as gr
import os
import json
import pandas as pd
from pathlib import Path
from PyPDF2 import PdfReader
from markdown2 import markdown
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from langchain_community.document_loaders import UnstructuredFileLoader  # Updated import

# ======================
# Core Extraction Logic
# ======================
ALLOWED_EXTENSIONS = {
    '.txt', '.md', '.json', '.csv', '.pdf',
    '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx'
}

def validate_file(filename):
    """Validate file existence and extension"""
    if not os.path.exists(filename):
        return False, "File does not exist."
    file_ext = os.path.splitext(filename)[1].lower()
    if file_ext not in ALLOWED_EXTENSIONS:
        return False, f"Unsupported file type '{file_ext}'"
    return True, ""

# File type specific extraction functions
def extract_text_from_txt(filename):
    with open(filename, 'r', encoding='utf-8') as f:
        return f.read()

def extract_text_from_md(filename):
    with open(filename, 'r', encoding='utf-8') as f:
        return markdown(f.read(), extras=['fenced-code-blocks'])

def extract_text_from_json(filename):
    with open(filename, 'r', encoding='utf-8') as f:
        return json.dumps(json.load(f), indent=2)

def extract_text_from_csv(filename):
    return pd.read_csv(filename).to_string()

def extract_text_from_pdf(filename):
    reader = PdfReader(filename)
    return "\n".join(page.extract_text() for page in reader.pages)

def extract_text_from_docx(filename):
    doc = Document(filename)
    return "\n".join(para.text for para in doc.paragraphs)

def extract_text_from_xlsx(filename):
    wb = load_workbook(filename)
    text = []
    for sheet in wb:
        for row in sheet.iter_rows(values_only=True):
            text.append(" ".join(str(cell) for cell in row if cell))
    return "\n".join(text)

def extract_text_from_pptx(filename):
    prs = Presentation(filename)
    return "\n".join(
        shape.text for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )

def extract_text_with_langchain(filename):
    return "\n".join(
        doc.page_content for doc in 
        UnstructuredFileLoader(filename).load()
    )

def extract_text(filename):
    """Main extraction router"""
    file_ext = os.path.splitext(filename)[1].lower()
    try:
        if file_ext == '.txt': return extract_text_from_txt(filename), ""
        elif file_ext == '.md': return extract_text_from_md(filename), ""
        elif file_ext == '.json': return extract_text_from_json(filename), ""
        elif file_ext == '.csv': return extract_text_from_csv(filename), ""
        elif file_ext == '.pdf': return extract_text_from_pdf(filename), ""
        elif file_ext in ('.docx', '.doc'): return extract_text_from_docx(filename), ""
        elif file_ext in ('.xlsx', '.xls'): return extract_text_from_xlsx(filename), ""
        elif file_ext in ('.pptx', '.ppt'): return extract_text_from_pptx(filename), ""
        else: return extract_text_with_langchain(filename), ""
    except Exception as e:
        return "", f"Extraction error: {str(e)}"

def save_extracted_text(filename, text):
    """Save extracted text to file"""
    try:
        output_path = f"{os.path.splitext(filename)[0]}_extracted.txt"
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        return f"Saved to {Path(output_path).name}"
    except Exception as e:
        return f"Save failed: {str(e)}"

# ======================
# Enhanced UI Components
# ======================
custom_css = """
.gradio-container {
    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Oxygen, Ubuntu, Cantarell;
    background: #f5f5f7 !important;
}
.block {
    background: white !important;
    border-radius: 18px !important;
    padding: 20px !important;
    margin: 10px 0 !important;
    box-shadow: 0 4px 6px rgba(0, 0, 0, 0.05) !important;
    border: 1px solid #e0e0e0 !important;
}
button {
    background: #007aff !important;
    color: white !important;
    border-radius: 12px !important;
    padding: 8px 16px !important;
    border: none !important;
    transition: all 0.2s ease;
}
button:hover {
    opacity: 0.9 !important;
    transform: scale(0.98);
}
.upload-button {
    background: #ffffff !important;
    border: 2px dashed #007aff !important;
    color: #007aff !important;
    padding: 2rem !important;
}
.preview-box {
    border-radius: 12px !important;
    border: 2px solid #e0e0e0 !important;
    padding: 15px !important;
}
"""

def create_ui():
    with gr.Blocks(theme=gr.themes.Soft(), css=custom_css) as demo:
        # Header Section
        gr.Markdown("# 📁 File Text Extractor")
        gr.Markdown("Extract text content from multiple files simultaneously. Supported formats: PDF, DOCX, XLSX, PPTX, TXT, and more.")

        # File Input Section
        with gr.Row():
            with gr.Column(scale=3):
                file_input = gr.File(
                    label="Select Files",
                    file_count="multiple",
                    file_types=list(ALLOWED_EXTENSIONS),
                    elem_classes="upload-button"
                )
            with gr.Column(scale=1):
                gr.Markdown("### Actions")
                with gr.Row():
                    extract_btn = gr.Button("Extract Text", variant="primary")
                    clear_btn = gr.Button("Clear All")

        # Processing Status
        status_box = gr.Markdown("## Status: Ready")

        # Preview Section with built-in copy button
        preview_box = gr.Textbox(
            label="Extracted Text Preview",
            interactive=True,
            lines=25,
            elem_classes="preview-box",
            show_copy_button=True  # Fixed copy functionality
        )

        # Footer
        gr.Markdown("---\n*Built with Gradio • iOS-inspired design • v2.0*")

        # ======================
        # Event Handling
        # ======================
        def process_files(files):
            outputs = []
            status = []
            total = len(files)
            
            for idx, file_info in enumerate(files, 1):
                file_path = file_info.name
                filename = Path(file_path).name
                base_msg = f"**Processing {idx}/{total}:** `{filename}`"
                
                # Validation
                valid, valid_msg = validate_file(file_path)
                if not valid:
                    status.append(f"{base_msg}\n❌ Validation failed: {valid_msg}")
                    continue
                
                # Extraction
                text, err = extract_text(file_path)
                if err:
                    status.append(f"{base_msg}\n❌ Extraction failed: {err}")
                    continue
                
                # Saving
                save_result = save_extracted_text(file_path, text)
                outputs.append(f"=== {filename} ===\n{text}\n")
                status.append(f"{base_msg}\n✅ Success\n📥 {save_result}")
            
            return {
                status_box: "\n\n".join(status),
                preview_box: "\n\n".join(outputs)
            }

        extract_btn.click(
            process_files,
            inputs=file_input,
            outputs=[status_box, preview_box]
        )

        clear_btn.click(
            lambda: [None, "## Status: Ready", ""],
            outputs=[file_input, status_box, preview_box]
        )

    return demo

# ======================
# Application Launch
# ======================
if __name__ == "__main__":
    ui = create_ui()
    ui.launch()
  

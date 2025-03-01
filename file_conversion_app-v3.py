import gradio as gr
import os
import json
import pandas as pd
import logging
import traceback
from pathlib import Path
from PyPDF2 import PdfReader
from markdown2 import markdown
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from langchain_community.document_loaders import UnstructuredFileLoader

# New imports for additional file types
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup
import xlrd

# Set up logging
logging.basicConfig(level=logging.INFO, 
                   format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# ======================
# Core Extraction Logic
# ======================
ALLOWED_EXTENSIONS = {
    '.txt', '.md', '.json', '.csv', '.pdf',
    '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx',
    '.epub'  # Added ePub support
}

def validate_file(filename):
    """Validate file existence and extension"""
    if not os.path.exists(filename):
        return False, "File does not exist."
    file_ext = os.path.splitext(filename)[1].lower()
    if file_ext not in ALLOWED_EXTENSIONS:
        return False, f"Unsupported file type '{file_ext}'. Supported types: {', '.join(ALLOWED_EXTENSIONS)}"
    return True, ""

# File type specific extraction functions
def extract_text_from_txt(filename):
    try:
        with open(filename, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        # Try different encodings if utf-8 fails
        for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
            try:
                with open(filename, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        raise Exception("Failed to decode text file with multiple encodings")

def extract_text_from_md(filename):
    with open(filename, 'r', encoding='utf-8') as f:
        return markdown(f.read(), extras=['fenced-code-blocks', 'tables', 'header-ids'])

def extract_text_from_json(filename):
    with open(filename, 'r', encoding='utf-8') as f:
        data = json.load(f)
        return json.dumps(data, indent=2, ensure_ascii=False)

def extract_text_from_csv(filename):
    try:
        df = pd.read_csv(filename)
        return df.to_string(index=False)
    except pd.errors.EmptyDataError:
        return "CSV file is empty"
    except Exception as e:
        # Try with different encodings and delimiters
        try:
            df = pd.read_csv(filename, encoding='latin-1')
            return df.to_string(index=False)
        except:
            try:
                df = pd.read_csv(filename, sep=';')
                return df.to_string(index=False)
            except:
                raise Exception(f"Failed to parse CSV: {str(e)}")

def extract_text_from_pdf(filename):
    reader = PdfReader(filename)
    text = []
    
    # Extract text from pages
    for i, page in enumerate(reader.pages):
        page_text = page.extract_text()
        if page_text:
            text.append(f"--- Page {i+1} ---\n{page_text}")
        else:
            text.append(f"--- Page {i+1} [No extractable text] ---")
    
    # Try to extract form fields if present
    try:
        fields = reader.get_fields()
        if fields:
            form_data = []
            for field, value in fields.items():
                if value:
                    form_data.append(f"{field}: {value}")
            if form_data:
                text.append("\n--- Form Data ---\n" + "\n".join(form_data))
    except:
        pass
    
    return "\n\n".join(text)

def extract_text_from_docx(filename):
    doc = Document(filename)
    text = []
    
    # Extract document properties if available
    try:
        core_props = doc.core_properties
        props = []
        if hasattr(core_props, 'title') and core_props.title:
            props.append(f"Title: {core_props.title}")
        if hasattr(core_props, 'author') and core_props.author:
            props.append(f"Author: {core_props.author}")
        if props:
            text.append("--- Document Properties ---\n" + "\n".join(props))
    except:
        pass
    
    # Extract paragraphs
    para_text = []
    for para in doc.paragraphs:
        if para.text.strip():
            para_text.append(para.text)
    
    if para_text:
        text.append("--- Content ---\n" + "\n".join(para_text))
    
    # Extract tables
    tables_text = []
    for i, table in enumerate(doc.tables):
        tables_text.append(f"--- Table {i+1} ---")
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            if row_text.strip():
                tables_text.append(row_text)
    
    if tables_text:
        text.append("\n".join(tables_text))
    
    return "\n\n".join(text)

def extract_text_from_xlsx(filename):
    wb = load_workbook(filename, data_only=True)  # data_only=True to get values instead of formulas
    text = []
    
    for sheet in wb:
        sheet_name = sheet.title
        text.append(f"\n--- Sheet: {sheet_name} ---\n")
        
        # Find the maximum column with data
        max_col = sheet.max_column
        max_row = sheet.max_row
        
        # Extract data rows
        for row in range(1, max_row + 1):
            row_values = []
            for col in range(1, max_col + 1):
                cell_value = sheet.cell(row=row, column=col).value
                row_values.append(str(cell_value) if cell_value is not None else "")
            
            if any(val.strip() for val in row_values):  # Only add if there's actual content
                text.append(" | ".join(row_values))
    
    return "\n".join(text)

def extract_text_from_xls(filename):
    """Extract text from legacy Excel .xls files"""
    workbook = xlrd.open_workbook(filename)
    text = []
    
    for sheet_idx in range(workbook.nsheets):
        sheet = workbook.sheet_by_index(sheet_idx)
        sheet_name = sheet.name
        text.append(f"\n--- Sheet: {sheet_name} ---\n")
        
        for row_idx in range(sheet.nrows):
            row_values = sheet.row_values(row_idx)
            row_text = " | ".join(str(cell) for cell in row_values if cell)
            if row_text.strip():
                text.append(row_text)
    
    return "\n".join(text)

def extract_text_from_pptx(filename):
    prs = Presentation(filename)
    text = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        slide_text.append(f"--- Slide {i+1} ---")
        
        # Extract from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        
        # Include notes if available
        if hasattr(slide, 'notes_slide') and slide.notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            slide_text.append(f"[Notes: {slide.notes_slide.notes_text_frame.text.strip()}]")
        
        text.append("\n".join(slide_text))
    
    return "\n\n".join(text)

def extract_text_from_epub(filename):
    """Extract text from EPUB e-books"""
    book = epub.read_epub(filename)
    text = []
    
    # Get metadata
    title = book.get_metadata('DC', 'title')
    creator = book.get_metadata('DC', 'creator')
    
    if title:
        text.append(f"Title: {title[0][0]}")
    if creator:
        text.append(f"Author: {creator[0][0]}")
    
    text.append("--- Content ---")
    
    # Extract content from HTML
    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            # Extract text from HTML content
            soup = BeautifulSoup(item.get_content(), 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.extract()
            
            content = soup.get_text(separator='\n')
            # Clean up whitespace
            content = '\n'.join(line.strip() for line in content.splitlines() if line.strip())
            if content:
                text.append(content)
    
    return "\n\n".join(text)

def extract_text_with_langchain(filename):
    """Use LangChain's UnstructuredFileLoader as a fallback"""
    try:
        documents = UnstructuredFileLoader(filename).load()
        return "\n".join(doc.page_content for doc in documents)
    except Exception as e:
        logger.error(f"LangChain extraction failed: {str(e)}")
        raise Exception(f"LangChain extraction failed: {str(e)}")

def extract_text(filename):
    """Main extraction router with improved error handling"""
    file_ext = os.path.splitext(filename)[1].lower()
    try:
        if file_ext == '.txt': return extract_text_from_txt(filename), ""
        elif file_ext == '.md': return extract_text_from_md(filename), ""
        elif file_ext == '.json': return extract_text_from_json(filename), ""
        elif file_ext == '.csv': return extract_text_from_csv(filename), ""
        elif file_ext == '.pdf': return extract_text_from_pdf(filename), ""
        elif file_ext == '.docx': return extract_text_from_docx(filename), ""
        elif file_ext == '.doc': 
            try:
                return extract_text_from_docx(filename), ""
            except:
                return extract_text_with_langchain(filename), "Note: Used UnstructuredFileLoader for .doc"
        elif file_ext == '.xlsx': return extract_text_from_xlsx(filename), ""
        elif file_ext == '.xls': return extract_text_from_xls(filename), ""
        elif file_ext in ('.pptx', '.ppt'): return extract_text_from_pptx(filename), ""
        elif file_ext == '.epub': return extract_text_from_epub(filename), ""
        else: 
            # Try with UnstructuredFileLoader as a fallback
            return extract_text_with_langchain(filename), f"Note: Used fallback extractor for {file_ext}"
    except Exception as e:
        logger.error(f"Extraction error for {filename}: {str(e)}")
        logger.error(traceback.format_exc())
        
        # Provide more specific error messages based on file type
        if file_ext == '.pdf':
            return "", f"PDF extraction error: {str(e)}. File might be encrypted, image-based, or damaged."
        elif file_ext in ('.doc', '.docx'):
            return "", f"Word document error: {str(e)}. File might be corrupted or password protected."
        elif file_ext in ('.xls', '.xlsx'):
            return "", f"Excel file error: {str(e)}. File might be corrupted or password protected."
        elif file_ext == '.epub':
            return "", f"EPUB error: {str(e)}. File might be corrupted or in an unsupported format."
        else:
            return "", f"Extraction error: {str(e)}"

def save_extracted_text(filename, text):
    """Save extracted text to file with better error handling"""
    try:
        # Create a clean output filename
        base_name = os.path.splitext(os.path.basename(filename))[0]
        output_path = f"{base_name}_extracted.txt"
        
        # Avoid overwriting existing files
        counter = 1
        while os.path.exists(output_path):
            output_path = f"{base_name}_extracted_{counter}.txt"
            counter += 1
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        return f"Saved to {Path(output_path).name}"
    except Exception as e:
        logger.error(f"Save error: {str(e)}")
        return f"Save failed: {str(e)}"

def save_all_text(text, output_filename=None):
    """Save all extracted text to a single file"""
    try:
        if not output_filename:
            output_filename = "extracted_text.txt"
            
            # Avoid overwriting existing files
            counter = 1
            while os.path.exists(output_filename):
                output_filename = f"extracted_text_{counter}.txt"
                counter += 1
        
        with open(output_filename, 'w', encoding='utf-8') as f:
            f.write(text)
        return f"Saved all text to {output_filename}"
    except Exception as e:
        logger.error(f"Save all text error: {str(e)}")
        return f"Save failed: {str(e)}"

# ======================
# Enhanced UI Components
# ======================
custom_css = """
.gradio-container {
    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Oxygen, Ubuntu, Cantarell;
    background: #f5f5f7 !important;
}
.block {
    background: white !important;
    border-radius: 18px !important;
    padding: 20px !important;
    margin: 10px 0 !important;
    box-shadow: 0 4px 6px rgba(0, 0, 0, 0.05) !important;
    border: 1px solid #e0e0e0 !important;
}
button {
    background: #007aff !important;
    color: white !important;
    border-radius: 12px !important;
    padding: 8px 16px !important;
    border: none !important;
    transition: all 0.2s ease;
}
button:hover {
    opacity: 0.9 !important;
    transform: scale(0.98);
}
.upload-button {
    background: #ffffff !important;
    border: 2px dashed #007aff !important;
    color: #007aff !important;
    padding: 2rem !important;
}
.preview-box {
    border-radius: 12px !important;
    border: 2px solid #e0e0e0 !important;
    padding: 15px !important;
}
.save-button {
    background: #34c759 !important;
}
"""

def create_ui():
    with gr.Blocks(theme=gr.themes.Soft(), css=custom_css) as demo:
        # Header Section
        gr.Markdown("# 📁 File Text Extractor")
        gr.Markdown("Extract text content from various file formats. Supported formats: PDF, DOCX, XLSX, PPTX, EPUB, TXT, and more.")

        # File Input Section
        with gr.Row():
            with gr.Column(scale=3):
                file_input = gr.File(
                    label="Select Files",
                    file_count="multiple",
                    file_types=list(ALLOWED_EXTENSIONS),
                    elem_classes="upload-button"
                )
            with gr.Column(scale=1):
                gr.Markdown("### Actions")
                with gr.Row():
                    extract_btn = gr.Button("Extract Text", variant="primary")
                    clear_btn = gr.Button("Clear All")

        # Processing Status
        status_box = gr.Markdown("## Status: Ready")

        # Preview Section with built-in copy button
        preview_box = gr.Textbox(
            label="Extracted Text Preview",
            interactive=True,
            lines=25,
            elem_classes="preview-box",
            show_copy_button=True
        )
        
        # Save Button (new addition)
        with gr.Row():
            save_btn = gr.Button("Save Extracted Text", variant="primary", elem_classes="save-button")
            save_filename = gr.Textbox(label="Save Filename (optional)", placeholder="extracted_text.txt")
        
        save_status = gr.Markdown("") # To show save status

        # Footer
        gr.Markdown("---\n*Built with Gradio • iOS-inspired design • v3.0*")

        # ======================
        # Event Handling
        # ======================
        def process_files(files):
            if not files:
                return {
                    status_box: "## Status: No files selected",
                    preview_box: ""
                }
                
            outputs = []
            status = []
            total = len(files)
            
            for idx, file_info in enumerate(files, 1):
                file_path = file_info.name
                filename = Path(file_path).name
                base_msg = f"**Processing {idx}/{total}:** `{filename}`"
                
                # Validation
                valid, valid_msg = validate_file(file_path)
                if not valid:
                    status.append(f"{base_msg}\n❌ Validation failed: {valid_msg}")
                    continue
                
                # Extraction
                try:
                    text, note = extract_text(file_path)
                    if not text:
                        status.append(f"{base_msg}\n❌ Extraction failed: No text content found")
                        continue
                        
                    outputs.append(f"=== {filename} ===\n{text}\n")
                    status_message = f"{base_msg}\n✅ Success"
                    if note:
                        status_message += f"\nℹ️ {note}"
                    status.append(status_message)
                except Exception as e:
                    status.append(f"{base_msg}\n❌ Extraction failed: {str(e)}")
            
            if not outputs:
                return {
                    status_box: "\n\n".join(status),
                    preview_box: "No text content could be extracted."
                }
                
            return {
                status_box: "\n\n".join(status),
                preview_box: "\n\n".join(outputs)
            }
        
        def save_text_content(text, filename):
            if not text:
                return "⚠️ Nothing to save - please extract text first"
            
            # Use custom filename if provided, otherwise use default
            custom_filename = filename.strip() if filename and filename.strip() else None
            result = save_all_text(text, custom_filename)
            return f"📄 {result}"

        extract_btn.click(
            process_files,
            inputs=file_input,
            outputs=[status_box, preview_box]
        )

        clear_btn.click(
            lambda: [None, "## Status: Ready", "", ""],
            outputs=[file_input, status_box, preview_box, save_status]
        )
        
        save_btn.click(
            save_text_content,
            inputs=[preview_box, save_filename],
            outputs=save_status
        )

    return demo

# ======================
# Application Launch
# ======================
if __name__ == "__main__":
    ui = create_ui()
    ui.launch()

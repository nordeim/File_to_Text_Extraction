# Install the necessary packages before running this script, for example:
# !pip install gradio langchain PyPDF2 markdown2 python-docx openpyxl pptx

import os
import gradio as gr
from PyPDF2 import PdfReader
import markdown2
import json
import csv
from docx import Document
import openpyxl
from pptx import Presentation
from langchain.document_loaders import UnstructuredFileLoader

def validate_file_type(file_path):
    """
    Check if the file extension is one of the allowed types.
    Allowed types:
      - Text-based: .txt, .md, .json, .csv, .pdf
      - Microsoft Office: .doc, .docx, .xls, .xlsx, .ppt, .pptx
    """
    allowed_extensions = {".txt", ".md", ".json", ".csv", ".pdf", 
                          ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx"}
    ext = os.path.splitext(file_path)[1].lower()
    return ext in allowed_extensions

def extract_text(file_path):
    """
    Extract text from the file based on its extension.
    Uses specialized libraries for each file type.
    Falls back to LangChain's UnstructuredFileLoader for unsupported types.
    """
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as file:
                return file.read()

        elif ext == ".md":
            with open(file_path, "r", encoding="utf-8") as file:
                # Convert Markdown to HTML
                return markdown2.markdown(file.read())

        elif ext == ".json":
            with open(file_path, "r", encoding="utf-8") as file:
                data = json.load(file)
            return json.dumps(data, indent=4)

        elif ext == ".csv":
            with open(file_path, "r", encoding="utf-8") as file:
                return file.read()

        elif ext == ".pdf":
            reader = PdfReader(file_path)
            text = ""
            # Extract text from each page of the PDF
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            return text

        elif ext in [".doc", ".docx"]:
            doc = Document(file_path)
            return "\n".join(paragraph.text for paragraph in doc.paragraphs)

        elif ext in [".xls", ".xlsx"]:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text = ""
            # Iterate through each sheet and row to form a text string
            for sheet in workbook:
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    text += row_text + "\n"
            return text

        elif ext in [".ppt", ".pptx"]:
            presentation = Presentation(file_path)
            text = ""
            # Iterate through each slide and extract text from shapes
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text

        else:
            # Use LangChain's UnstructuredFileLoader as a fallback solution
            loader = UnstructuredFileLoader(file_path)
            documents = loader.load()
            return "\n".join([doc.page_content for doc in documents])

    except Exception as e:
        return f"Error extracting text: {str(e)}"

def process_file(file):
    """
    Gradio callback function to process the uploaded file.
    - Validates the file type.
    - Extracts the text.
    - Saves the extracted text to a new file (original filename appended with '_extracted.txt').
    - Returns a message including the location of the saved output and the text content.
    """
    temp_filepath = file.name
    
    # Validate file type
    if not validate_file_type(temp_filepath):
        return f"File type not supported: {os.path.splitext(temp_filepath)[1]}"
    
    # Extract text using the appropriate method based on file type
    extracted_text = extract_text(temp_filepath)
    
    # Create output filename by appending '_extracted.txt' to the original base name
    base, _ = os.path.splitext(temp_filepath)
    output_filepath = base + "_extracted.txt"
    
    try:
        with open(output_filepath, "w", encoding="utf-8") as output_file:
            output_file.write(extracted_text)
        response = f"Extracted text saved to: {output_filepath}\n\nExtracted Content:\n{extracted_text}"
    except Exception as e:
        response = f"Extraction succeeded but error saving file: {str(e)}\n\nExtracted Content:\n{extracted_text}"
    
    return response

def create_ui():
    """
    Create and launch the Gradio web interface.
    The interface displays:
      - A file upload widget.
      - A textbox to show the extraction result.
    """
    interface = gr.Interface(
        fn=process_file,
        inputs=gr.File(label="Upload a file"),
        outputs=gr.Textbox(label="Output"),
        title="File to Text Converter",
        description=("Upload a file to extract its text content. "
                     "Supported formats: .txt, .md, .json, .csv, .pdf, .doc, .docx, .xls, .xlsx, .ppt, .pptx.")
    )
    interface.launch()

if __name__ == "__main__":
    create_ui()

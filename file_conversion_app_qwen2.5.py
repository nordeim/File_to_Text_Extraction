import os
import json
import pandas as pd
from PyPDF2 import PdfReader
import markdown2
from docx import Document
from pptx import Presentation
from langchain.document_loaders import UnstructuredFileLoader
import gradio as gr

# Define allowed file extensions
ALLOWED_EXTENSIONS = ['.txt', '.md', '.json', '.csv', '.pdf', '.docx', '.xlsx', '.pptx']

def extract_text_from_file(file_path):
    """
    Extracts text from the given file based on its extension.
    Falls back to LangChain if no specialized library is available.
    """
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    try:
        if ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        elif ext == '.md':
            with open(file_path, 'r', encoding='utf-8') as f:
                return markdown2.markdown(f.read())
        elif ext == '.json':
            with open(file_path, 'r', encoding='utf-8') as f:
                return json.dumps(json.load(f), indent=4)
        elif ext == '.csv':
            df = pd.read_csv(file_path)
            return df.to_string(index=False)
        elif ext == '.pdf':
            reader = PdfReader(file_path)
            return ''.join(page.extract_text() or '' for page in reader.pages)
        elif ext == '.docx':
            doc = Document(file_path)
            return '\n'.join(paragraph.text for paragraph in doc.paragraphs)
        elif ext == '.xlsx':
            df = pd.read_excel(file_path)
            return df.to_string(index=False)
        elif ext == '.pptx':
            presentation = Presentation(file_path)
            return '\n'.join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))
        else:
            # Fallback to LangChain for unsupported types
            loader = UnstructuredFileLoader(file_path)
            return loader.load()[0].page_content
    except Exception as e:
        raise ValueError(f"Error processing file: {e}")

def process_file(file_name):
    """
    Main function to validate file type, extract content, and save as .txt.
    """
    # Validate file extension
    _, ext = os.path.splitext(file_name)
    ext = ext.lower()
    if ext not in ALLOWED_EXTENSIONS:
        return f"Unsupported file type: {ext}. Allowed types: {', '.join(ALLOWED_EXTENSIONS)}"

    try:
        # Extract text
        extracted_text = extract_text_from_file(file_name)

        # Save output to a new file
        output_file = f"{os.path.splitext(file_name)[0]}_output.txt"
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(extracted_text)

        return f"Text extracted successfully! Saved to: {output_file}"
    except Exception as e:
        return f"An error occurred: {str(e)}"

# Gradio Interface
with gr.Blocks() as demo:
    gr.Markdown("# File to Text Converter")
    file_input = gr.Textbox(label="Enter File Path")
    output = gr.Textbox(label="Result")
    submit_button = gr.Button("Convert")

    submit_button.click(process_file, inputs=file_input, outputs=output)

# Launch the app
demo.launch()

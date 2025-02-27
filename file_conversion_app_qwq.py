import gradio as gr
import os
import json
import pandas as pd
from PyPDF2 import PdfReader
from markdown2 import markdown
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from langchain_community.document_loaders import UnstructuredFileLoader

ALLOWED_EXTENSIONS = {
    '.txt', '.md', '.json', '.csv', '.pdf',
    '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx'
}

def validate_file(filename):
    """Validates the file extension and checks if the file exists."""
    if not os.path.exists(filename):
        return False, "Error: File does not exist."
    
    file_ext = os.path.splitext(filename)[1].lower()
    if file_ext not in ALLOWED_EXTENSIONS:
        return False, f"Error: Unsupported file type '{file_ext}'. Allowed types: {', '.join(ALLOWED_EXTENSIONS)}"
    
    return True, ""

def extract_text_from_txt(filename):
    """Extracts text from a plain text file."""
    with open(filename, 'r', encoding='utf-8') as f:
        return f.read()

def extract_text_from_md(filename):
    """Extracts text from a Markdown file."""
    with open(filename, 'r', encoding='utf-8') as f:
        md_content = f.read()
        return markdown(md_content, extras=['fenced-code-blocks'])

def extract_text_from_json(filename):
    """Extracts text from a JSON file."""
    with open(filename, 'r', encoding='utf-8') as f:
        json_content = json.load(f)
        return json.dumps(json_content, indent=2)

def extract_text_from_csv(filename):
    """Extracts text from a CSV file."""
    df = pd.read_csv(filename)
    return df.to_string()

def extract_text_from_pdf(filename):
    """Extracts text from a PDF file."""
    reader = PdfReader(filename)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def extract_text_from_docx(filename):
    """Extracts text from a DOCX file."""
    doc = Document(filename)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_xlsx(filename):
    """Extracts text from an XLSX file."""
    wb = load_workbook(filename)
    text = ""
    for sheet in wb:
        for row in sheet.rows:
            for cell in row:
                if cell.value:
                    text += str(cell.value) + " "
            text += "\n"
    return text

def extract_text_from_pptx(filename):
    """Extracts text from a PPTX file."""
    prs = Presentation(filename)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_with_langchain(filename):
    """Fallback: Extracts text using LangChain for unsupported file types."""
    loader = UnstructuredFileLoader(filename)
    documents = loader.load()
    return "\n".join([doc.page_content for doc in documents])

def extract_text(filename):
    """
    Extracts text from the given file based on its extension.
    
    Args:
        filename (str): Path to the input file.
    
    Returns:
        tuple: (str, str) - (extracted_text, error_message)
    """
    file_ext = os.path.splitext(filename)[1].lower()
    
    try:
        if file_ext == '.txt':
            return extract_text_from_txt(filename), ""
        elif file_ext == '.md':
            return extract_text_from_md(filename), ""
        elif file_ext == '.json':
            return extract_text_from_json(filename), ""
        elif file_ext == '.csv':
            return extract_text_from_csv(filename), ""
        elif file_ext == '.pdf':
            return extract_text_from_pdf(filename), ""
        elif file_ext in ('.docx', '.doc'):
            return extract_text_from_docx(filename), ""
        elif file_ext in ('.xlsx', '.xls'):
            return extract_text_from_xlsx(filename), ""
        elif file_ext in ('.pptx', '.ppt'):
            return extract_text_from_pptx(filename), ""
        else:
            return extract_text_with_langchain(filename), ""
    except Exception as e:
        return "", f"Error extracting text: {str(e)}"

def save_extracted_text(original_filename, text):
    """Saves extracted text with original filename prefix"""
    base_name = os.path.splitext(os.path.basename(original_filename))[0]
    output_filename = f"{base_name}_extracted.txt"
    try:
        with open(output_filename, 'w', encoding='utf-8') as f:
            f.write(text)
        return f"Saved to: {output_filename}"
    except Exception as e:
        return f"Save error: {str(e)}"

def process_files(files):
    """Process multiple files and return combined results"""
    all_results = []
    for file in files:
        temp_path = file.name
        original_name = os.path.basename(temp_path)
        
        valid, msg = validate_file(temp_path)
        if not valid:
            all_results.append(f"❌ {original_name}: {msg}")
            continue
        
        text, error = extract_text(temp_path)
        if error:
            all_results.append(f"❌ {original_name}: {error}")
            continue
        
        save_msg = save_extracted_text(original_name, text)
        all_results.append(f"✅ {original_name}\n{save_msg}\n\n{text}")
    
    return "\n".join(all_results)

css = """
.header { text-align: center; padding: 20px; background: #f0f0f0; border-radius: 10px; }
.file-uploader { margin: 20px 0; }
.process-btn { background: #4CAF50; color: white; padding: 12px 24px; border-radius: 8px; }
.process-btn:hover { background: #45a049; }
.output-box { border: 1px solid #ddd; border-radius: 8px; padding: 15px; height: 350px; }
.copy-btn { margin: 10px 0; padding: 8px 16px; background: #2196F3; color: white; border-radius: 5px; }
"""

def create_ui():
    """Create enhanced Gradio interface"""
    with gr.Blocks(css=css) as demo:
        with gr.Row():
            gr.Markdown("# 📄 File Text Extractor", elem_classes="header")
        
        with gr.Row():
            file_upload = gr.File(
                label="Upload Files",
                file_count="multiple",
                elem_classes="file-uploader"
            )
        
        with gr.Row():
            submit_btn = gr.Button("Extract Text", elem_classes="process-btn")
            copy_btn = gr.Button("📋 Copy Text", elem_classes="copy-btn")
        
        with gr.Row():
            output_text = gr.Textbox(
                label="Extracted Text Preview",
                lines=15,
                interactive=False,
                elem_classes="output-box",
                elem_id="output_text"
            )
        
        submit_btn.click(
            fn=process_files,
            inputs=file_upload,
            outputs=output_text
        )
        
        copy_btn.click(
            fn=None,
            js="() => { navigator.clipboard.writeText(document.getElementById('output_text').value); }"
        )
    
    return demo

if __name__ == "__main__":
    ui = create_ui()
    ui.launch()

